"""
ingest_pptx.py — Convert PPTX files to Markdown using python-pptx.

Part of process-ai ingest pipeline. Extracts slide content as structured
markdown: slide titles become H2 headings, text shapes become paragraphs,
speaker notes become blockquotes, and embedded images are extracted.
Group shapes (e.g. SmartArt) are noted as placeholders for future vision
pipeline (v2).

Contract: emits JSON to stdout — { ok, format, slides, markdown, images, metadata }
"""

import argparse
import os
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ingest_common import (
    emit_error,
    emit_success,
    extract_title_from_md,
    slugify_for_path,
    save_image_bytes,
    ensure_images_dir,
)


def _shape_has_text(shape) -> bool:
    """Check if a shape has a text frame with content."""
    if not shape.has_text_frame:
        return False
    return bool(shape.text_frame.text.strip())


def _shape_text(shape) -> str:
    """Get cleaned text from a shape's text frame."""
    if not shape.has_text_frame:
        return ''
    return shape.text_frame.text.strip()


def _extract_images_from_slide(slide, output_dir, slug, slide_num) -> list:
    """Extract images from a slide. Returns list of (rel_path, description)."""
    images = []
    img_dir = ensure_images_dir(output_dir, slug)
    img_count = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                ext = image.content_type.split('/')[-1] if '/' in image.content_type else 'png'
                if ext == 'jpeg':
                    ext = 'jpg'
                data = image.blob

                img_path = save_image_bytes(data, ext, img_dir,
                                            slide_num * 100 + img_count)
                img_count += 1

                # Try to get the image description from alt text
                desc = shape.alt_text if hasattr(shape, 'alt_text') else ''
                images.append((img_path, desc))
            except Exception:
                # Image extraction failure — skip silently
                pass

    return images


def _is_bullet(text: str) -> bool:
    """Heuristic: detect if text looks like a bullet item."""
    stripped = text.strip()
    return bool(re.match(r'^[•\-–—*·▪▸►]\s', stripped))


def _group_shape_description(shape) -> str:
    """Produce a placeholder description for group shapes (e.g. SmartArt)."""
    shape_name = shape.name or 'Diagrama'
    shape_type = str(shape.shape_type) if shape.shape_type else 'GROUP'
    count = (len(shape.shapes) if hasattr(shape, 'shapes') and
             shape.shapes else 0)
    return (f"<!-- [DIAGRAM: {shape_name} ({shape_type}, "
            f"{count} sub-shapes)] -->")


def convert_pptx(pptx_path: str, output_dir: str):
    """Convert a PPTX file to Markdown. Returns (md_path, images_rel, metadata)."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        raise RuntimeError(f"Não foi possível abrir o arquivo PPTX (formato inválido?): {e}")

    filename = os.path.basename(pptx_path)
    base_name = os.path.splitext(filename)[0]
    slug = slugify_for_path(base_name)

    md_name = f"{slug}.md"
    md_path = os.path.join(output_dir, md_name)

    # Extract metadata from core properties
    cp = prs.core_properties
    metadata = {
        'source_file': filename,
        'title': cp.title or base_name,
        'author': cp.author or '',
        'created': cp.created.isoformat() if cp.created else '',
        'modified': cp.modified.isoformat() if cp.modified else '',
    }

    all_images_rel = []
    lines_out = []

    # Use the presentation title as H1 if available
    if cp.title:
        lines_out.append(f'# {cp.title}\n')

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines_out.append(f'---\n')
        lines_out.append(f'## Slide {slide_num}\n')

        has_title = False

        for shape in slide.shapes:
            # Handle GroupShapes (SmartArt, grouped diagrams)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                lines_out.append(_group_shape_description(shape))
                lines_out.append('')
                continue

            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_images = _extract_images_from_slide(slide, output_dir, slug, slide_num)
                # Filter images we haven't seen
                for img_path, desc in slide_images:
                    img_rel = os.path.relpath(img_path, output_dir).replace('\\', '/')
                    if img_rel not in all_images_rel:
                        all_images_rel.append(img_rel)
                        lines_out.append('<figure>')
                        lines_out.append(f'  <img src="{img_rel}" alt="{desc or "Slide image"}" />')
                        if desc:
                            lines_out.append(f'  <figcaption>{desc}</figcaption>')
                        else:
                            lines_out.append('  <figcaption><!-- [IMAGE] --></figcaption>')
                        lines_out.append('</figure>\n')
                continue

            # Placeholder (title, subtitle, body)
            if shape.is_placeholder:
                ph = shape.placeholder_format
                ph_type = ph.type if ph.type is not None else None
                # CENTER_TITLE=3, TITLE=1, SUBTITLE=2
                is_title_ph = ph_type is not None and ph_type in (1, 3)  # TITLE or CENTER_TITLE

                if _shape_has_text(shape):
                    text = _shape_text(shape)
                    if is_title_ph and not has_title:
                        lines_out.append(f'### {text}\n')
                        has_title = True
                    elif is_title_ph:
                        # Additional title — treat as body
                        lines_out.append(text + '\n')
                    else:
                        # Subtitle or body
                        for line in text.split('\n'):
                            stripped = line.strip()
                            if not stripped:
                                continue
                            if _is_bullet(stripped):
                                lines_out.append(f'- {stripped[1:].strip()}\n')
                            else:
                                lines_out.append(stripped + '\n')
                continue

            # Regular text shapes
            if shape.has_text_frame and _shape_text(shape):
                text = _shape_text(shape)
                for line in text.split('\n'):
                    stripped = line.strip()
                    if not stripped:
                        continue
                    if _is_bullet(stripped):
                        lines_out.append(f'- {stripped[1:].strip()}\n')
                    else:
                        lines_out.append(stripped + '\n')
                continue

        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide
            notes_text = notes.notes_text_frame.text.strip() if notes.notes_text_frame else ''
            if notes_text:
                lines_out.append('> **Notas do apresentador:** ' + notes_text.replace('\n', ' ') + '\n')

    md_body = '\n'.join(lines_out)
    # Collapse excessive blank lines
    md_body = re.sub(r'\n{3,}', '\n\n', md_body)

    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(md_body)

    # Extract title from markdown
    extracted_title = extract_title_from_md(md_body)
    if extracted_title:
        metadata['title'] = extracted_title

    slide_count = len(prs.slides)

    return md_path, all_images_rel, metadata, slide_count


def main():
    parser = argparse.ArgumentParser(
        description='Convert PPTX to Markdown (process-ai ingest)',
    )
    parser.add_argument('--input', required=True, help='Path to PPTX file')
    parser.add_argument('--output-dir', required=True,
                        help='Directory for output markdown + images')
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.input)
    output_dir = os.path.abspath(args.output_dir)

    if not os.path.isfile(pptx_path):
        emit_error(f"Arquivo não encontrado: {pptx_path}")
        return

    try:
        os.makedirs(output_dir, exist_ok=True)
        md_path, images, metadata, slide_count = convert_pptx(pptx_path, output_dir)
        md_rel = os.path.relpath(md_path, output_dir).replace('\\', '/')

        emit_success(
            format='pptx',
            markdown_rel=md_rel,
            images=images,
            metadata=metadata,
            slides=slide_count,
        )
    except Exception as e:
        emit_error(f"Falha ao converter PPTX '{os.path.basename(pptx_path)}': {str(e)}")


if __name__ == "__main__":
    main()
